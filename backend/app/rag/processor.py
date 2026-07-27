"""Main RAG processor - orchestrates everything."""
from pathlib import Path
import mimetypes
from typing import Dict, List, Optional
import os
import json
import uuid

# Replace Qdrant with LanceDB service
from app.services.vector_store import (
    add_chunks, 
    search_chunks, 
    get_all_chunks, 
    get_base_dir
)
from app.database import get_db

from .ocr import OCRProcessor
from .hybrid_search import HybridSearcher
from .chunking import chunk_text, add_context_to_chunks, generate_summary

# Use FastEmbed for reranking (ONNX) instead of Transformers/Torch
try:
    from fastembed import TextRerank
    RERANK_MODEL = "BAAI/bge-reranker-base"
except ImportError:
    TextRerank = None

class RAGProcessor:
    def __init__(self, qdrant_path: str = None):
        # qdrant_path is ignored now, but kept for compatibility if passed
        print("Initializing RAG models (LanceDB + FastEmbed)...")
        
        # Initialize OCR (EasyOCR wrapper)
        try:
            self.ocr = OCRProcessor()
            print("OCR processor initialized")
        except Exception as e:
            print(f"Warning: OCR processor failed to initialize: {e}")
            self.ocr = None
        
        # Initialize Reranker (FastEmbed)
        try:
            if TextRerank:
                self.reranker = TextRerank(model_name=RERANK_MODEL)
                print("Reranker initialized (FastEmbed)")
            else:
                self.reranker = None
                print("Warning: fastembed not installed, reranking disabled")
        except Exception as e:
            print(f"Warning: Reranker failed to initialize: {e}")
            self.reranker = None
        
        # Initialize Hybrid Searcher (BM25)
        try:
            self.hybrid = HybridSearcher()
            print("Hybrid searcher initialized")
        except Exception as e:
            print(f"Warning: Hybrid searcher failed to initialize: {e}")
            self.hybrid = None
        
        print("RAG Processor initialized.")
    
    def process_file(self, file_path: str, conversation_id: str, user_id: int = 1) -> Dict:
        """Process any file type and store in LanceDB and SQLite."""
        # Ensure absolute path for file:// URI compatibility
        file_path = Path(file_path).resolve()
        mime_type, _ = mimetypes.guess_type(file_path)
        
        print(f"Processing file: {file_path}, type: {mime_type}")
        
        try:
            # Route to appropriate processor
            text = ""
            method = "text"
            
            if mime_type and mime_type.startswith('image/'):
                text = self._process_image(file_path)
                method = "vision"
            elif file_path.suffix.lower() == '.pdf':
                text = self._process_pdf(file_path)
            elif file_path.suffix.lower() == '.docx':
                text = self._process_docx(file_path)
            elif file_path.suffix.lower() == '.pptx':
                text = self._process_pptx(file_path)
            elif file_path.suffix.lower() in ['.xlsx', '.xls']:
                text = self._process_excel(file_path)
            elif file_path.suffix.lower() in ['.txt', '.md', '.py', '.js', '.ts', '.html', '.css', '.json']:
                text = file_path.read_text(encoding='utf-8', errors='ignore')
            else:
                return {"status": "error", "message": f"Unsupported file type: {file_path.suffix}"}
            
            if not text.strip():
                return {"status": "error", "message": "No text extracted from file"}
            
            # Generate summary for context
            doc_summary = generate_summary(text)
            
            # Chunk with context
            chunks = chunk_text(text, chunk_size=500, overlap=50)
            contextual_chunks = add_context_to_chunks(chunks, doc_summary, file_path.name)
            
            print(f"Storing {len(contextual_chunks)} chunks for {file_path.name} in LanceDB and SQLite...")
            
            # 1. Track in SQLite
            doc_id = str(uuid.uuid4())
            try:
                db = get_db()
                db.execute(
                    "INSERT INTO documents (id, user_id, filename, file_type, chunk_count) VALUES (?, ?, ?, ?, ?)",
                    (doc_id, user_id, file_path.name, mime_type or file_path.suffix, len(contextual_chunks))
                )
                
                # Batch insert vector chunks bridge
                chunk_records = []
                for i, chunk_text in enumerate(contextual_chunks):
                    chunk_id = f"{conversation_id}_{file_path.name}_{i}"
                    chunk_records.append((chunk_id, doc_id, i, chunk_text))
                
                db.executemany(
                    "INSERT INTO vector_chunks (id, document_id, chunk_index, chunk_text) VALUES (?, ?, ?, ?)",
                    chunk_records
                )
                db.commit()
                db.close()
            except Exception as e:
                print(f"Warning: Failed to track document in SQLite: {e}")

            # 2. Store in LanceDB
            db_chunks = []
            for i, chunk_text in enumerate(contextual_chunks):
                # Construct metadata
                metadata = {
                    "conversation_id": conversation_id,
                    "file_name": file_path.name,
                    "file_type": mime_type or file_path.suffix,
                    "chunk_index": i,
                    "processing_method": method,
                    "parent_summary": doc_summary,
                    "sqlite_doc_id": doc_id
                }
                
                db_chunks.append({
                    "id": f"{conversation_id}_{file_path.name}_{i}",
                    "document_id": conversation_id, 
                    "chunk_text": chunk_text,
                    "metadata": json.dumps(metadata)
                })
            
            add_chunks(db_chunks)
            
            # Index for BM25 (update index for this conversation)
            self._update_bm25_index(conversation_id)
            
            print(f"Successfully processed and stored {file_path.name}")
            
            return {
                "status": "success",
                "chunks": len(chunks),
                "file_name": file_path.name,
                "processing_method": method,
                "extracted_text": text
            }
            
        except Exception as e:
            import traceback
            traceback.print_exc()
            return {"status": "error", "message": str(e)}
    
    def _process_image(self, file_path: Path) -> str:
        """Process image using EasyOCR only."""
        image_path_str = str(file_path)
        print(f"Processing image: {image_path_str}")
        
        # Only use OCR
        if self.ocr is not None:
            try:
                print(f"Attempting OCR extraction for {file_path.name}...")
                ocr_text = self.ocr.extract_text(image_path_str)
                print(f"OCR extraction completed, extracted {len(ocr_text)} characters")
                
                if ocr_text and ocr_text.strip():
                    return ocr_text
                else:
                    print(f"OCR extraction produced empty text for {file_path.name}")
            except Exception as e:
                print(f"OCR extraction failed for {file_path.name}: {e}")
        else:
            print("OCR processor not initialized")
        
        # If OCR fails or is empty, return minimal description
        print(f"Warning: OCR failed or produced no text for image: {file_path.name}")
        
        # Try to provide at least some context based on filename
        filename = file_path.name.lower()
        description_parts = []
        
        if any(ext in filename for ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']):
            description_parts.append("This is an image file")
        
        # Add basic description based on filename patterns
        if 'screenshot' in filename:
            description_parts.append("It appears to be a screenshot")
        elif 'photo' in filename or 'picture' in filename:
            description_parts.append("It appears to be a photograph")
        elif 'diagram' in filename or 'chart' in filename:
            description_parts.append("It appears to be a diagram or chart")
        elif 'document' in filename or 'scan' in filename:
            description_parts.append("It appears to be a scanned document")
        
        # Add file size info if available
        try:
            file_size = file_path.stat().st_size
            if file_size > 0:
                description_parts.append(f"File size: {file_size:,} bytes")
        except:
            pass
        
        if description_parts:
            return f"[Image Processing Note: Could not extract meaningful text from image. {', '.join(description_parts)}.]"
        else:
            return f"[Image Processing Note: Could not extract meaningful text from image {file_path.name}. The image may not contain readable text or the processing models failed.]"
    
    def _process_pdf(self, file_path: Path) -> str:
        """Extract text and OCR images from PDF using pdfplumber."""
        import pdfplumber
        import tempfile
        
        try:
            text_content = []
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = [f"### Page {i+1}"]
                    
                    # 1. Extract standard text
                    content = page.extract_text()
                    if content:
                        page_text.append(content)
                    
                    # 2. Extract tables (pdfplumber specialty)
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            table_str = "\n| " + " | ".join([str(c).replace("\n", " ") if c else "" for c in table[0]]) + " |\n"
                            table_str += "| " + " | ".join(["---" for _ in table[0]]) + " |\n"
                            for row in table[1:]:
                                table_str += "| " + " | ".join([str(c).replace("\n", " ") if c else "" for c in row]) + " |\n"
                            page_text.append(table_str)
                    
                    # 3. OCR Images on page
                    # If very little text was extracted, try OCR on the whole page
                    if not content or len(content.strip()) < 50:
                        try:
                            # Render page to image for OCR
                            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                                page.to_image(resolution=200).save(tmp.name)
                                tmp_path = Path(tmp.name)
                            
                            ocr_result = self._process_image(tmp_path)
                            if ocr_result and "Image Processing Note" not in ocr_result and ocr_result.strip():
                                page_text.append(f"\n[OCR Text from Page Image:\n{ocr_result}\n]")
                            
                            if os.path.exists(tmp_path):
                                os.unlink(tmp_path)
                        except Exception as e:
                            print(f"Failed to OCR PDF page {i+1}: {e}")

                    text_content.append("\n".join(page_text))
            
            return "\n\n---\n\n".join(text_content)
        except Exception as e:
            print(f"PDF processing failed: {e}")
            # Fallback to simple pypdf if pdfplumber fails
            try:
                import pypdf
                with open(file_path, 'rb') as f:
                    reader = pypdf.PdfReader(f)
                    return "\n\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
            except:
                return f"Error processing PDF: {str(e)}"
    
    def _process_docx(self, file_path: Path) -> str:
        """Extract text from DOCX."""
        import docx
        doc = docx.Document(file_path)
        return "\n\n".join([para.text for para in doc.paragraphs])

    def _process_pptx(self, file_path: Path) -> str:
        """Robustly extract text from PPTX."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import tempfile
        
        def get_shape_text(shape):
            text = ""
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = "".join(run.text for run in paragraph.runs).strip()
                    if para_text:
                        text += para_text + "\n"
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    text += "| " + " | ".join([cell.text_frame.text.replace("\n", " ").strip() for cell in row.cells]) + " |\n"
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for s in shape.shapes:
                    text += get_shape_text(s) + "\n"
            elif hasattr(shape, "text") and shape.text:
                text += shape.text.strip() + "\n"
            return text.strip()

        try:
            prs = Presentation(file_path)
            all_slides_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_parts = [f"--- SLIDE {i+1} ---"]
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        slide_parts.append(f"TITLE: {title_text}")
                
                for shape in slide.shapes:
                    if shape == slide.shapes.title: continue
                    shape_text = get_shape_text(shape)
                    if shape_text:
                        slide_parts.append(shape_text)
                    
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            with tempfile.NamedTemporaryFile(suffix=f".{image.ext}", delete=False) as tmp:
                                tmp.write(image.blob)
                                tmp_path = Path(tmp.name)
                            
                            ocr_text = self._process_image(tmp_path)
                            if ocr_text and "[Image Processing Note:" not in ocr_text and ocr_text.strip():
                                slide_parts.append(f"[TEXT DETECTED IN IMAGE: {ocr_text.strip()}]")
                            
                            if os.path.exists(tmp_path): os.unlink(tmp_path)
                        except Exception: pass

                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_parts.append(f"SPEAKER NOTES: {notes}")

                all_slides_content.append("\n".join(slide_parts))

            return "\n\n".join(all_slides_content)
        except Exception as e:
            return f"[Error processing PPTX {file_path.name}: {str(e)}]"
    
    def _process_excel(self, file_path: Path) -> str:
        """Extract text from Excel."""
        text = []
        try:
            # Try pylightxl
            if file_path.suffix.lower() == '.xlsx':
                try:
                    import pylightxl as xl
                    db = xl.readxl(fn=str(file_path))
                    for sheet in db.ws_names:
                        text.append(f"### Sheet: {sheet}")
                        rows = list(db.ws(ws=sheet).rows)
                        if not rows: continue
                        for r_idx, row in enumerate(rows, start=1):
                            row_parts = []
                            for c_idx, cell in enumerate(row):
                                if cell and str(cell).strip():
                                    row_parts.append(str(cell).replace("\n", " ").strip())
                            if row_parts:
                                text.append(f"Row {r_idx}: " + " | ".join(row_parts))
                        text.append("\n")
                    return "\n".join(text)
                except ImportError: pass
                
            # Fallback to openpyxl
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            for sheet in wb.worksheets:
                text.append(f"### Sheet: {sheet.title}")
                for r_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    row_parts = [str(c).replace("\n", " ").strip() for c in row if c]
                    if row_parts:
                        text.append(f"Row {r_idx}: " + " | ".join(row_parts))
                text.append("\n")
            return "\n".join(text)
        except Exception as e:
            return f"[Excel Processing Error: {str(e)}]"
    
    def _update_bm25_index(self, conversation_id: str):
        """Update BM25 index with conversation documents from LanceDB."""
        if self.hybrid is None:
            return
            
        try:
            # Retrieve all chunks for this conversation (document_id = conversation_id)
            chunks = get_all_chunks(conversation_id)
            documents = [c["chunk_text"] for c in chunks]
            
            if documents:
                self.hybrid.index_corpus(documents)
        except Exception as e:
            print(f"Failed to update BM25 index: {e}")
    
    def retrieve_context(self, query: str, conversation_id: str, top_k: int = 3) -> str:
        """Retrieve relevant context using hybrid search + reranking."""
        try:
            self._update_bm25_index(conversation_id)
            
            # Dense retrieval (LanceDB)
            dense_results = search_chunks(query, limit=15, document_id=conversation_id)
            dense_docs = [r["chunk_text"] for r in dense_results]
            
            # Sparse retrieval (BM25)
            sparse_docs = []
            if self.hybrid is not None:
                try:
                    sparse_docs = self.hybrid.search(query, top_k=15)
                except Exception as e:
                    print(f"Sparse retrieval failed: {e}")
            
            # Merge results
            merged = []
            if dense_docs and sparse_docs and self.hybrid is not None:
                try:
                    merged = self.hybrid.merge_results(dense_docs, sparse_docs)
                except Exception:
                    merged = dense_docs[:15]
            elif dense_docs:
                merged = dense_docs[:15]
            elif sparse_docs:
                merged = sparse_docs[:15]
            
            # Rerank
            if merged and self.reranker is not None:
                try:
                    # FastEmbed reranker expects list of strings
                    reranked = self.reranker.rank(query, merged[:20])
                    # reranked is iterator of Dict or similar? fastembed doc says:
                    # list of result objects/dicts. 
                    # Actually, TextRerank.rank returns list of {index, score, text} or similar.
                    # Let's check fastembed API usage.
                    # It returns generator of RerankResult(document=..., score=..., index=...)
                    
                    sorted_results = sorted(reranked, key=lambda x: x.score, reverse=True)
                    top_results = sorted_results[:top_k]
                    context = "\n\n---\n\n".join([r.document for r in top_results])
                    return context
                except Exception as e:
                    print(f"Reranking failed: {e}")
                    context = "\n\n---\n\n".join(merged[:top_k])
                    return context
            elif merged:
                context = "\n\n---\n\n".join(merged[:top_k])
                return context
            
            return ""
        except Exception as e:
            print(f"Retrieval failed: {e}")
            return ""
    
    def get_file_content(self, file_names: List[str], conversation_id: str) -> str:
        """
        Retrieve full content of specific files from LanceDB.
        """
        try:
            if not file_names:
                return ""
            
            print(f"Retrieving content for files: {file_names} in conv: {conversation_id}")
            
            # 1. Try to get from Vector Store
            found_files = {}
            # Get all chunks for the conversation and filter in memory
            all_chunks = get_all_chunks(conversation_id)
            
            for chunk in all_chunks:
                try:
                    meta = json.loads(chunk.get("metadata", "{}"))
                    fname = meta.get("file_name", "")
                    if fname in file_names:
                        if fname not in found_files:
                            found_files[fname] = []
                        found_files[fname].append((meta.get("chunk_index", 0), chunk.get("chunk_text", "")))
                except:
                    continue
            
            # 2. Process each requested file (same logic as before, just adapted source)
            context_parts = []
            files_with_errors = []
            files_not_found = []
            files_empty = []
            
            for fname in file_names:
                content = ""
                source = "vector_db"
                has_error = False
                
                if fname in found_files:
                    chunks = sorted(found_files[fname], key=lambda x: x[0])
                    content = "".join([c[1] for c in chunks])
                    if not content.strip():
                        has_error = True
                        files_empty.append(fname)
                else:
                    # Fallback to uploads folder (same logic as before)
                    source = "upload"
                    media_dir = os.getenv("MEDIA_DIR")
                    base_dir = get_base_dir()
                    
                    if media_dir:
                        if not os.path.isabs(media_dir):
                            uploads_dir = base_dir / media_dir / "uploads"
                        else:
                            uploads_dir = Path(media_dir) / "uploads"
                    else:
                        uploads_dir = base_dir / "media" / "uploads"
                    
                    file_path = uploads_dir / fname
                    
                    if file_path.exists():
                        try:
                            # Re-use process logic (without storing) or just read text
                            # For simplicity here, assume text or try basic read
                            content = file_path.read_text(errors='ignore')
                        except Exception:
                            files_not_found.append(fname)
                            continue
                    else:
                        files_not_found.append(fname)
                        continue
                
                if content:
                    context_parts.append(f"\n--- Content of Attached File: {fname} ({source}) ---\n{content}")
            
            return "\n".join(context_parts) if context_parts else "No content found."
            
        except Exception as e:
            print(f"Failed to get file content: {e}")
            return f"Error retrieving content: {e}"

    def search(self, query: str, conversation_id: str, top_k: int = 10) -> List[Dict]:
        """Structured search returning list of results with metadata."""
        try:
            results = search_chunks(query, limit=top_k, document_id=conversation_id)
            formatted_results = []
            for r in results:
                meta = json.loads(r.get("metadata", "{}"))
                formatted_results.append({
                    "content": r["chunk_text"],
                    "metadata": meta,
                    "similarity": r.get("score", 0.0) # LanceDB returns _distance, need to check if score available
                })
            return formatted_results
        except Exception as e:
            print(f"Search failed: {e}")
            return []
